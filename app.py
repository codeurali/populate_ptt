import collections.abc
import sys
from pptx import Presentation

import json
from flask import Flask, request, jsonify, render_template


def passToPTT(var_text):
  prs = Presentation('test.pptx')

  # Get slide shapes
  slide = prs.slides[1]
  shapes = slide.shapes

  # Get all placeholders id and text
  for shape in shapes:
    if not shape.has_text_frame:
      continue
    id = shape.placeholder_format.idx
    if id == 2: shape.text = var_text

  prs.save('test.pptx')


app = Flask(__name__)

@app.route('/')
def index():
  return render_template('index.html')




@app.route('/test', methods=['POST'])
def test():
  output = request.get_json()
  print(output)
  print(type(output))
  result = json.loads(output)
  print(result)
  print(type(result))
  passToPTT(result['arg1'])
  return result


if __name__ == '__main__':
  app.run(debug=True)